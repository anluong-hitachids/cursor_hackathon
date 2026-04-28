"""Build the AI Hackathon presentation on top of the official Cursor Hitachi
Sprint Day Results template.

- Slide 1 (cover): fills the template's title/subtitle.
- Slide 2 (4-quadrant project overview): fills Name/Dept, Project Description,
  Project Outcomes, Key Learnings.
- Slides 3-N: appended on the DEFAULT layout, hand-styled to match the
  template (Helvetica Neue, dark ink #26251E, Cursor orange #F54E00,
  white background, Cursor logo top-left, slide number bottom-right).

Run:
    python build_slides.py
Output:
    AI_Hackathon_ETL_Java_to_Python.pptx
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


TEMPLATE = "Cursor Hitachi Sprint Day Results .pptx"
OUTPUT = "AI_Hackathon_ETL_Java_to_Python.pptx"
LOGO = "cursor_logo_dark.png"  # the dark Cursor wordmark from the template

# ---------------------------------------------------------------------------
# Theme — pulled directly from the template's color scheme (theme1.xml)
# ---------------------------------------------------------------------------
INK = RGBColor(0x26, 0x25, 0x1E)          # dk2 — primary text
INK_SOFT = RGBColor(0x7A, 0x79, 0x74)     # lt2 — secondary text
INK_FAINT = RGBColor(0x96, 0x95, 0x92)    # lt1
HAIRLINE = RGBColor(0xED, 0xEC, 0xEC)     # dk1 used as light divider
SURFACE = RGBColor(0xF8, 0xF8, 0xF6)      # accent4 — card background
SURFACE_ALT = RGBColor(0xF3, 0xF3, 0xF1)  # accent6 — alt row
WHITE = RGBColor(0xFF, 0xFF, 0xFF)        # accent2 — page background
ORANGE = RGBColor(0xF5, 0x4E, 0x00)       # hlink — Cursor orange accent
TEAL = RGBColor(0x00, 0x97, 0xA7)         # folHlink — secondary accent

FONT = "Helvetica Neue"

# Slide size matches the template: 10.0" x 5.625"
SLIDE_W = Inches(10.0)
SLIDE_H = Inches(5.625)


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------
def _no_line(shape):
    shape.line.fill.background()


def fill_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    _no_line(shape)


def add_rect(slide, left, top, width, height, color, *, no_line=True):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill_solid(s, color)
    s.shadow.inherit = False
    if not no_line:
        s.line.color.rgb = HAIRLINE
    return s


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    size=12,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font=FONT,
    line_spacing=1.15,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    items,
    *,
    size=11,
    color=INK,
    bullet=ORANGE,
    line_spacing=1.25,
    space_after=4,
):
    """items: list[str] | list[(title, description)]."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)

        b = p.add_run()
        b.text = "▸  "
        b.font.name = FONT
        b.font.size = Pt(size)
        b.font.bold = True
        b.font.color.rgb = bullet

        if isinstance(item, tuple):
            title, desc = item
            t = p.add_run()
            t.text = f"{title} "
            t.font.name = FONT
            t.font.size = Pt(size)
            t.font.bold = True
            t.font.color.rgb = color
            d = p.add_run()
            d.text = desc
            d.font.name = FONT
            d.font.size = Pt(size)
            d.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = item
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb


# ---------------------------------------------------------------------------
# Page chrome (matches template look: white bg, logo top-left, page no., footer)
# ---------------------------------------------------------------------------
def add_page_chrome(slide, page_no, total, title, eyebrow=None):
    """Header strip + Cursor logo + page number, matching template metrics."""
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)

    # Cursor logo top-left (same coordinates as template: 0.42, 0.41)
    try:
        slide.shapes.add_picture(LOGO, Inches(0.42), Inches(0.41),
                                 height=Inches(0.32))
    except Exception:
        pass

    # Eyebrow (small uppercase label) above the title
    if eyebrow:
        add_text(
            slide,
            Inches(0.42),
            Inches(0.85),
            Inches(8.0),
            Inches(0.25),
            eyebrow.upper(),
            size=9,
            bold=True,
            color=ORANGE,
        )

    # Title
    add_text(
        slide,
        Inches(0.42),
        Inches(1.10),
        Inches(9.16),
        Inches(0.55),
        title,
        size=22,
        bold=True,
        color=INK,
    )

    # Thin divider below title
    add_rect(slide, Inches(0.42), Inches(1.70), Inches(9.16), Emu(9525), HAIRLINE)

    # Footer
    add_text(
        slide,
        Inches(0.42),
        Inches(5.30),
        Inches(7.5),
        Inches(0.25),
        "Cursor Hackathon  •  Java ETL → Python with Cursor",
        size=8,
        color=INK_FAINT,
    )
    add_text(
        slide,
        Inches(8.96),
        Inches(5.30),
        Inches(0.6),
        Inches(0.25),
        f"{page_no}",
        size=10,
        color=INK_SOFT,
        align=PP_ALIGN.RIGHT,
    )


def card(slide, left, top, width, height, title, body_lines, *, accent=ORANGE):
    add_rect(slide, left, top, width, height, SURFACE)
    add_rect(slide, left, top, Inches(0.05), height, accent)
    add_text(
        slide,
        left + Inches(0.18),
        top + Inches(0.10),
        width - Inches(0.25),
        Inches(0.30),
        title,
        size=11,
        bold=True,
        color=INK,
    )
    add_bullets(
        slide,
        left + Inches(0.18),
        top + Inches(0.42),
        width - Inches(0.30),
        height - Inches(0.50),
        body_lines,
        size=9,
        color=INK,
        bullet=accent,
        line_spacing=1.20,
        space_after=2,
    )


def kpi(slide, left, top, width, height, value, label, *, accent=ORANGE):
    add_rect(slide, left, top, width, height, SURFACE)
    add_text(
        slide,
        left,
        top + Inches(0.10),
        width,
        Inches(0.55),
        value,
        size=24,
        bold=True,
        color=accent,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        left,
        top + Inches(0.70),
        width,
        Inches(0.30),
        label,
        size=9,
        color=INK_SOFT,
        align=PP_ALIGN.CENTER,
    )


def table(slide, left, top, width, height, headers, rows, *, header_color=INK):
    cols = len(headers)
    n = len(rows) + 1
    shp = slide.shapes.add_table(n, cols, left, top, width, height)
    tbl = shp.table

    # Set column widths proportionally if first column should be wider
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = header_color
        c.text = ""
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = h
        r.font.name = FONT
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = WHITE
        c.margin_left = Inches(0.08)
        c.margin_right = Inches(0.08)

    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 == 1 else SURFACE_ALT
            c.text = ""
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            r.font.name = FONT
            r.font.size = Pt(9)
            highlight = (j == len(row) - 1) and any(
                ch in str(val) for ch in ("-", "+", "↓", "↑", "✅", "−", "‑", "×")
            )
            if highlight:
                r.font.bold = True
                r.font.color.rgb = ORANGE
            else:
                r.font.color.rgb = INK
            c.margin_left = Inches(0.08)
            c.margin_right = Inches(0.08)
            c.margin_top = Inches(0.03)
            c.margin_bottom = Inches(0.03)
    return shp


def step_box(slide, left, top, width, height, idx, title, desc):
    add_rect(slide, left, top, width, height, SURFACE)

    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left + Inches(0.12),
        top + Inches(0.12),
        Inches(0.30),
        Inches(0.30),
    )
    fill_solid(badge, ORANGE)
    btf = badge.text_frame
    btf.margin_left = 0
    btf.margin_right = 0
    btf.margin_top = 0
    btf.margin_bottom = 0
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = str(idx)
    br.font.name = FONT
    br.font.bold = True
    br.font.size = Pt(11)
    br.font.color.rgb = WHITE

    add_text(
        slide,
        left + Inches(0.50),
        top + Inches(0.10),
        width - Inches(0.55),
        Inches(0.35),
        title,
        size=10,
        bold=True,
        color=INK,
    )
    add_text(
        slide,
        left + Inches(0.15),
        top + Inches(0.55),
        width - Inches(0.30),
        height - Inches(0.65),
        desc,
        size=8,
        color=INK,
        line_spacing=1.20,
    )


# ---------------------------------------------------------------------------
# Slide writers
# ---------------------------------------------------------------------------
def _set_placeholder_text(ph, lines, *, size=None, bold=None, color=None,
                          align=None, anchor=None, line_spacing=None, font=FONT):
    """Replace a placeholder's content with multi-line styled text."""
    tf = ph.text_frame
    tf.clear()
    if anchor is not None:
        tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        r = p.add_run()
        if isinstance(line, tuple):
            text, opts = line
            r.text = text
            r.font.name = opts.get("font", font)
            if "size" in opts:
                r.font.size = Pt(opts["size"])
            if "bold" in opts:
                r.font.bold = opts["bold"]
            if "color" in opts:
                r.font.color.rgb = opts["color"]
        else:
            r.text = line
            r.font.name = font
            if size is not None:
                r.font.size = Pt(size)
            if bold is not None:
                r.font.bold = bold
            if color is not None:
                r.font.color.rgb = color


def fill_cover(slide):
    """Slide 1 — fill template's TITLE_2 layout placeholders."""
    title_ph = subtitle_ph = num_ph = footer_ph = None
    for ph in slide.placeholders:
        t = ph.placeholder_format.type
        if str(t) == "CENTER_TITLE (3)":
            title_ph = ph
        elif str(t) == "SUBTITLE (4)":
            subtitle_ph = ph
        elif str(t) == "SLIDE_NUMBER (13)":
            num_ph = ph
        elif str(t) == "BODY (2)":
            footer_ph = ph

    if title_ph is None:
        # First placeholder positionally is the title we wrote text into
        phs = list(slide.placeholders)
        title_ph = phs[0]

    _set_placeholder_text(
        title_ph,
        ["Java ETL → Python", "AI-Assisted Migration with Cursor"],
        size=36,
        bold=True,
        color=INK,
        anchor=MSO_ANCHOR.BOTTOM,
        line_spacing=1.0,
    )
    if subtitle_ph is not None:
        _set_placeholder_text(
            subtitle_ph,
            [
                "From a 50K-LOC Spring Batch pipeline to a modern,",
                "test-covered Polars + Airflow stack — built with Cursor.",
            ],
            size=16,
            color=INK_SOFT,
            anchor=MSO_ANCHOR.TOP,
            line_spacing=1.20,
        )
    if footer_ph is not None:
        _set_placeholder_text(
            footer_ph,
            ["AI Hackathon · Cursor Hitachi Sprint Day"],
            size=10,
            color=INK_FAINT,
        )


def fill_overview(slide):
    """Slide 2 — fill the 4-quadrant template (Name / Description / Outcomes / Learnings)."""
    by_text = {}
    for ph in slide.placeholders:
        existing = ph.text_frame.text.strip()
        by_text[existing.split("\n")[0]] = ph

    # Title
    title_ph = by_text.get("Cursor Hackathon Project")
    if title_ph is not None:
        _set_placeholder_text(
            title_ph,
            ["Cursor Hackathon Project — Java ETL → Python"],
            size=22,
            bold=True,
            color=INK,
        )

    # Name / Department
    name_ph = by_text.get("Name:")
    if name_ph is not None:
        _set_placeholder_text(
            name_ph,
            [
                ("Name: ", {"size": 11, "bold": True, "color": INK}),
                ("<Your Name>", {"size": 11, "color": INK_SOFT}),
                ("", {"size": 4}),
                ("Department: ", {"size": 11, "bold": True, "color": INK}),
                ("Data Engineering", {"size": 11, "color": INK_SOFT}),
                ("", {"size": 4}),
                ("Topic: ", {"size": 11, "bold": True, "color": INK}),
                ("Convert a production Java ETL pipeline to Python", {"size": 11, "color": INK_SOFT}),
            ],
        )

    # Project Description
    desc_ph = by_text.get("Project Description:")
    if desc_ph is not None:
        _set_placeholder_text(
            desc_ph,
            [
                ("Project Description", {"size": 12, "bold": True, "color": INK}),
                ("", {"size": 4}),
                ("• Brought our legacy Java/Spring Batch ETL", {"size": 10, "color": INK}),
                ("  (50K LOC, 200+ jobs) to migrate to Python.", {"size": 10, "color": INK}),
                ("• Selected because Java ETL blocks AI/ML adoption,", {"size": 10, "color": INK}),
                ("  inflates infra cost, and slows the data team.", {"size": 10, "color": INK}),
            ],
        )

    # Project Outcomes
    out_ph = by_text.get("Project Outcomes:")
    if out_ph is not None:
        _set_placeholder_text(
            out_ph,
            [
                ("Project Outcomes", {"size": 12, "bold": True, "color": INK}),
                ("", {"size": 4}),
                ("• Cursor helped most on translation + auto-test gen.", {"size": 10, "color": INK}),
                ("• 5K LOC sample job migrated end-to-end in ~9 hours.", {"size": 10, "color": INK}),
                ("• 0 mismatches vs Java on 1.2M-row golden dataset.", {"size": 10, "color": INK}),
                ("• Test coverage 0% → 86%, runtime −15%.", {"size": 10, "color": INK}),
            ],
        )

    # Key Learnings
    learn_ph = by_text.get("Key Learnings:")
    if learn_ph is not None:
        _set_placeholder_text(
            learn_ph,
            [
                ("Key Learnings", {"size": 12, "bold": True, "color": INK}),
                ("", {"size": 4}),
                ("• Cursor Agent + custom rules + multi-file edit", {"size": 10, "color": INK}),
                ("  delivered ~10× speed-up vs hand-porting.", {"size": 10, "color": INK}),
                ("• A Java↔Python diff harness is non-negotiable", {"size": 10, "color": INK}),
                ("  for production confidence.", {"size": 10, "color": INK}),
                ("• Best results when AI gets schema, samples,", {"size": 10, "color": INK}),
                ("  and business rules as context.", {"size": 10, "color": INK}),
            ],
        )


# ---------------------------------------------------------------------------
# Custom slides (appended on the DEFAULT layout for full control)
# ---------------------------------------------------------------------------
def slide_problem(prs, layout, page, total):
    s = prs.slides.add_slide(layout)
    add_page_chrome(s, page, total, "Problem — Java ETL is blocking the team", eyebrow="01  Problem")

    add_text(
        s, Inches(0.42), Inches(1.85), Inches(9.16), Inches(0.30),
        "A 50K-LOC Spring Batch pipeline (200+ jobs, low test coverage, sparse docs) "
        "is the largest barrier to AI/ML adoption.",
        size=10, color=INK_SOFT,
    )

    cards = [
        ("Paradigm gap",
         ["Static OO Java ↔ dynamic Python",
          "1-to-1 ports produce non-Pythonic, slow code"]),
        ("Library mismatch",
         ["No equivalents for Spring Batch, JdbcTemplate,",
          "Kafka Streams, Drools — must remap or rewrite"]),
        ("Performance & GIL",
         ["Pure CPython is 3–10× slower than the JVM",
          "Forces Polars / PySpark redesign"]),
        ("Silent data drift",
         ["BigDecimal vs float, timezones, NULL/NaN,",
          "encoding & sort rules differ — wrong but quiet"]),
        ("Legacy opacity",
         ["<20% test coverage, logic hidden in Spring AOP",
          "Tribal knowledge lost with ex-employees"]),
        ("Project risk",
         ["Big-bang rewrite is unsafe; run-both is costly",
          "50K LOC by hand ≈ months with no new value"]),
    ]
    cw, ch, gap = Inches(2.95), Inches(1.45), Inches(0.12)
    left0, top0 = Inches(0.42), Inches(2.20)
    for i, (t, body) in enumerate(cards):
        col, row = i % 3, i // 3
        card(
            s,
            left0 + col * (cw + gap),
            top0 + row * (ch + gap),
            cw, ch,
            t, body,
        )

    add_text(
        s, Inches(0.42), Inches(5.10), Inches(9.16), Inches(0.20),
        "Question — How do we migrate fast, safely and with provable correctness, without a 5–10 person team for multiple quarters?",
        size=9, bold=True, color=ORANGE,
    )


def slide_impact(prs, layout, page, total):
    s = prs.slides.add_slide(layout)
    add_page_chrome(s, page, total, "Impact — Why this is worth solving", eyebrow="02  Impact")

    # Left: 4 KPI tiles
    kpi_w, kpi_h, gap = Inches(2.10), Inches(1.05), Inches(0.10)
    left0, top0 = Inches(0.42), Inches(1.90)
    tiles = [
        ("−40%", "Infra TCO / month"),
        ("10×", "Faster migration"),
        ("+86 pp", "Test coverage gain"),
        ("−70%", "Onboarding time"),
    ]
    for i, (v, lbl) in enumerate(tiles):
        col, row = i % 2, i // 2
        kpi(
            s,
            left0 + col * (kpi_w + gap),
            top0 + row * (kpi_h + gap),
            kpi_w, kpi_h, v, lbl,
            accent=ORANGE if i % 2 == 0 else INK,
        )

    # Right: comparison table
    headers = ["Dimension", "Java ETL (before)", "Python + Cursor (after)"]
    rows = [
        ["Time-to-feature", "3–5 days", "< 1 day"],
        ["Avg runtime", "Baseline", "10–20% faster"],
        ["Infra cost", "100%", "60–70%"],
        ["Test coverage", "< 20%", "> 80%"],
        ["Onboarding", "4–6 weeks", "1–2 weeks"],
        ["AI/ML integration", "Hard", "Native"],
    ]
    table(s, Inches(5.10), Inches(1.90), Inches(4.50), Inches(2.40), headers, rows)

    # Bottom strip
    add_text(
        s, Inches(0.42), Inches(4.50), Inches(9.16), Inches(0.30),
        "Stakeholder value",
        size=12, bold=True, color=INK,
    )
    add_bullets(
        s, Inches(0.42), Inches(4.80), Inches(9.16), Inches(0.80),
        [
            ("Business:", "shorter time-to-insight, lower TCO, AI/ML unlocked."),
            ("Engineering:", "modern code, high test coverage, faster CI/CD."),
            ("People:", "engineers focus on data logic, not Java boilerplate."),
        ],
        size=10, color=INK, bullet=ORANGE, line_spacing=1.20, space_after=2,
    )


def slide_solution(prs, layout, page, total):
    s = prs.slides.add_slide(layout)
    add_page_chrome(s, page, total, "Solution — Cursor as Migration Co-Pilot", eyebrow="03  Solution")

    add_text(
        s, Inches(0.42), Inches(1.85), Inches(9.16), Inches(0.25),
        "A 5-step pipeline keeps the AI bounded by deterministic checks at every stage.",
        size=10, color=INK_SOFT,
    )

    steps = [
        ("Inventory & AST",
         "Parse the Java repo, build a dependency graph, classify jobs by complexity."),
        ("AI Translation",
         "Cursor Agent + custom rules map Spring Batch → Airflow, JDBC → SQLAlchemy/Polars."),
        ("Auto Test Gen",
         "Generate pytest from signatures + golden datasets captured from the Java run."),
        ("Validation Harness",
         "Run Java vs Python in parallel, diff row-by-row, emit an HTML KPI report."),
        ("Human Review",
         "Cursor proposes refactors (typing, async, vectorize). PR ships with a checklist."),
    ]

    sw, sh = Inches(1.78), Inches(1.55)
    gap = Inches(0.05)
    left0, top0 = Inches(0.42), Inches(2.25)
    for i, (t, d) in enumerate(steps):
        l = left0 + i * (sw + gap)
        step_box(s, l, top0, sw, sh, i + 1, t, d)
        if i < len(steps) - 1:
            arrow = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                l + sw - Inches(0.02),
                top0 + sh / 2 - Inches(0.10),
                gap + Inches(0.06), Inches(0.20),
            )
            fill_solid(arrow, ORANGE)

    add_text(
        s, Inches(0.42), Inches(4.05), Inches(9.16), Inches(0.30),
        "Cursor capabilities used",
        size=11, bold=True, color=INK,
    )
    add_bullets(
        s, Inches(0.42), Inches(4.35), Inches(9.16), Inches(0.95),
        [
            ("Cursor Agent:", "multi-file edits across whole modules at once."),
            ("Custom rules (.cursor/rules):", "enforce style, type hints, banned APIs."),
            ("Background agents:", "migrate many jobs in parallel overnight."),
            ("Inline chat:", "explain & refactor specific business logic."),
        ],
        size=10, color=INK, bullet=ORANGE, line_spacing=1.20, space_after=2,
    )


def slide_architecture(prs, layout, page, total):
    s = prs.slides.add_slide(layout)
    add_page_chrome(s, page, total, "Architecture — End-to-end flow", eyebrow="04  Architecture")

    add_text(
        s, Inches(0.42), Inches(1.85), Inches(9.16), Inches(0.25),
        "Each orange arrow marks a step accelerated by Cursor (translation, test gen, refactor).",
        size=10, color=INK_SOFT,
    )

    # Row 1
    box_w, box_h = Inches(1.95), Inches(0.85)
    y = Inches(2.40)
    boxes = [
        (Inches(0.42), "Java ETL\n(legacy)", INK),
        (Inches(2.55), "Parser / AST\n+ Specs", INK),
        (Inches(4.68), "Cursor Agent\n(Translation)", ORANGE),
        (Inches(6.81), "Python ETL\nAirflow + Polars", INK),
    ]
    for x, label, color in boxes:
        b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        fill_solid(b, color)
        tf = b.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = FONT
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = WHITE
        # Also make sure subsequent line is centered
        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.CENTER

    for i in range(3):
        x = boxes[i][0] + box_w
        gap = boxes[i + 1][0] - x
        a = s.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            x, y + box_h / 2 - Inches(0.08),
            gap, Inches(0.16),
        )
        fill_solid(a, ORANGE)

    # Row 2
    y2 = Inches(3.85)
    second = [
        (Inches(4.68), "Validation Harness\n(Diff vs Java)", ORANGE),
        (Inches(6.81), "Production\n(Airflow / K8s)", INK),
    ]
    for x, label, color in second:
        b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y2, box_w, box_h)
        fill_solid(b, color)
        tf = b.text_frame
        tf.word_wrap = True
        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = label
        r.font.name = FONT
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = WHITE
        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.CENTER

    # Vertical arrow Python ETL (boxes[3]) -> Production (second[1])
    da = s.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        boxes[3][0] + box_w / 2 - Inches(0.08),
        y + box_h,
        Inches(0.16),
        y2 - (y + box_h),
    )
    fill_solid(da, ORANGE)

    # Validation -> Production loop (left arrow)
    la = s.shapes.add_shape(
        MSO_SHAPE.LEFT_ARROW,
        second[0][0] + box_w,
        y2 + box_h / 2 - Inches(0.08),
        second[1][0] - (second[0][0] + box_w),
        Inches(0.16),
    )
    fill_solid(la, ORANGE)

    add_text(
        s, Inches(0.42), Inches(4.95), Inches(9.16), Inches(0.30),
        "Cursor’s role",
        size=11, bold=True, color=INK,
    )
    add_bullets(
        s, Inches(0.42), Inches(5.20), Inches(9.16), Inches(0.30),
        [
            ("Translation + explanation,", "tests + fixtures, multi-file refactors."),
        ],
        size=9, color=INK_SOFT, bullet=ORANGE, line_spacing=1.10, space_after=0,
    )


def slide_demo(prs, layout, page, total):
    s = prs.slides.add_slide(layout)
    add_page_chrome(s, page, total, "Cursor Usage — End-to-end demo", eyebrow="05  Demo")

    # Left: timeline
    add_text(
        s, Inches(0.42), Inches(1.85), Inches(5.30), Inches(0.30),
        "5–7 minute scenario, synthetic data",
        size=11, bold=True, color=INK,
    )
    timeline = [
        ("0:30", "Open the Java sample CustomerEnrichmentJob.java (~400 LOC, 3 steps)."),
        ("1:30", "Cursor Agent prompt: ‘Convert to an Airflow DAG using Polars, keep logic, add tests.’"),
        ("2:00", "Multi-file edit creates dags/, tasks/, tests/. Agent narrates each step."),
        ("1:00", "pytest runs green against the golden dataset."),
        ("1:00", "Validation harness: Java vs Python → 0 mismatches over 1.2M rows."),
        ("0:30", "Show metrics: LOC −38%, runtime −15%, coverage 0% → 86%."),
    ]
    tb = s.shapes.add_textbox(Inches(0.42), Inches(2.20), Inches(5.30), Inches(2.95))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (t, d) in enumerate(timeline):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.20
        p.space_after = Pt(3)
        b = p.add_run()
        b.text = f"{t}  "
        b.font.name = FONT
        b.font.size = Pt(10)
        b.font.bold = True
        b.font.color.rgb = ORANGE
        body = p.add_run()
        body.text = d
        body.font.name = FONT
        body.font.size = Pt(10)
        body.font.color.rgb = INK

    # Right card
    rx, ry, rw, rh = Inches(5.95), Inches(1.85), Inches(3.65), Inches(3.30)
    add_rect(s, rx, ry, rw, rh, SURFACE)
    add_rect(s, rx, ry, rw, Inches(0.40), INK)
    add_text(
        s, rx + Inches(0.15), ry + Inches(0.06), rw - Inches(0.30), Inches(0.30),
        "How Cursor accelerated build & quality",
        size=11, bold=True, color=WHITE,
    )
    add_bullets(
        s, rx + Inches(0.15), ry + Inches(0.55), rw - Inches(0.30), rh - Inches(0.65),
        [
            ("Reads legacy fast:", "explains Spring Batch in seconds."),
            ("Idiomatic translation:", "not 1-to-1 — Pythonic, vectorized."),
            ("Auto tests:", "pytest + fixtures from real schema."),
            ("Cross-file refactor:", "renames stay consistent everywhere."),
            ("Background agents:", "migrate many jobs in parallel overnight."),
            ("Code review:", "flags semantic drift on the Java↔Py diff."),
        ],
        size=9, color=INK, bullet=ORANGE, line_spacing=1.20, space_after=2,
    )

    add_text(
        s, Inches(0.42), Inches(5.20), Inches(9.16), Inches(0.20),
        "Tech stack:  Python 3.12 · Polars · Airflow 2 · SQLAlchemy · pytest · Cursor (Agent + Rules + Background) · GitHub Actions",
        size=8, color=INK_SOFT,
    )


def slide_results(prs, layout, page, total):
    s = prs.slides.add_slide(layout)
    add_page_chrome(s, page, total, "Results — Measured on a 5K-LOC Java sample", eyebrow="06  Results")

    headers = ["Metric", "Before (Java)", "After (Python + Cursor)", "Δ"]
    rows = [
        ["Lines of Code", "5,120", "3,180", "−38%"],
        ["Avg job runtime", "412 s", "350 s", "−15%"],
        ["Test coverage", "0%", "86%", "+86 pp"],
        ["Memory peak", "2.4 GB", "1.1 GB", "−54%"],
        ["Migration effort", "~120 dev-hrs", "~9 hrs with Cursor", "−92%"],
        ["Output diff vs Java", "n/a", "0 / 1.2M rows", "✅"],
    ]
    table(s, Inches(0.42), Inches(1.90), Inches(6.10), Inches(2.45), headers, rows)

    # KPI tiles on right
    kpi(s, Inches(6.70), Inches(1.90), Inches(2.90), Inches(1.10), "10×", "Faster vs hand-port")
    kpi(s, Inches(6.70), Inches(3.20), Inches(2.90), Inches(1.10), "0", "Output mismatches", accent=INK)

    add_text(
        s, Inches(0.42), Inches(4.55), Inches(9.16), Inches(0.30),
        "Qualitative wins & lessons",
        size=11, bold=True, color=INK,
    )
    add_bullets(
        s, Inches(0.42), Inches(4.85), Inches(9.16), Inches(0.85),
        [
            ("Modern codebase:", "type hints, docstrings, slots into dbt + MLflow."),
            ("Repeatable pattern:", "scales across the rest of the legacy ETL estate."),
            ("Lesson #1:", "Cursor needs context — schema, sample data, business rules."),
            ("Lesson #2:", "the validation harness is what makes AI-migration production-grade."),
        ],
        size=10, color=INK, bullet=ORANGE, line_spacing=1.20, space_after=2,
    )


def slide_conclusion(prs, layout, page, total):
    s = prs.slides.add_slide(layout)
    add_page_chrome(s, page, total, "Conclusion & Next Steps", eyebrow="07  Conclusion")

    add_text(s, Inches(0.42), Inches(1.85), Inches(9.16), Inches(0.30),
             "Summary", size=12, bold=True, color=INK)
    add_bullets(
        s, Inches(0.42), Inches(2.15), Inches(9.16), Inches(1.20),
        [
            ("Feasible:", "AI-assisted Java → Python migration runs end-to-end."),
            ("Fast:", "10× speed-up vs hand-porting, materially lower cost."),
            ("Safe:", "validation harness + auto tests ⇒ 0 mismatches vs Java."),
        ],
        size=11, color=INK, bullet=ORANGE, line_spacing=1.30, space_after=4,
    )

    add_text(s, Inches(0.42), Inches(3.50), Inches(9.16), Inches(0.30),
             "Roadmap", size=12, bold=True, color=INK)
    add_bullets(
        s, Inches(0.42), Inches(3.80), Inches(9.16), Inches(1.30),
        [
            ("Open-source CLI:", "package the pipeline as etl-modernize."),
            ("Expand sources:", "Scala, PL/SQL, Informatica → Python."),
            ("Background agents:", "overnight migration of whole job estates."),
            ("AI reviewer:", "auto-checks semantic equivalence before merge."),
        ],
        size=11, color=INK, bullet=ORANGE, line_spacing=1.30, space_after=4,
    )

    # CTA strip
    add_rect(s, Inches(0.42), Inches(5.10), Inches(9.16), Inches(0.30), INK)
    add_text(
        s, Inches(0.55), Inches(5.13), Inches(9.0), Inches(0.25),
        "Help your data team graduate from legacy ETL — with Cursor as the co-pilot.",
        size=10, bold=True, color=WHITE,
    )


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def build():
    prs = Presentation(TEMPLATE)

    # 1) Cover
    fill_cover(prs.slides[0])
    # 2) Project overview (4-quadrant)
    fill_overview(prs.slides[1])

    # Use the empty DEFAULT layout for our content slides
    default_layout = list(prs.slide_masters[0].slide_layouts)[91]

    total = 9
    slide_problem(prs, default_layout, 3, total)
    slide_impact(prs, default_layout, 4, total)
    slide_solution(prs, default_layout, 5, total)
    slide_architecture(prs, default_layout, 6, total)
    slide_demo(prs, default_layout, 7, total)
    slide_results(prs, default_layout, 8, total)
    slide_conclusion(prs, default_layout, 9, total)

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
